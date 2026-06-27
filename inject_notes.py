#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
把 FY26_EO_Training_备注讲解.md 里的逐页备注，一次性写进本地的 PPTX「备注」框。

为什么需要这个脚本：
  原始 pptx 体积 245MB，无法在云端环境里来回搬运回写。
  在你自己电脑上跑这个脚本，就能把全部 63 页备注直接注入你本地的原始文件，
  不用上传、不受体积限制、几秒钟搞定。

用法（在 Mac/Windows 终端里）：
  1) 安装依赖（只需一次）：
       pip3 install python-pptx
  2) 把本脚本和 FY26_EO_Training_备注讲解.md 放到与 pptx 同一目录，然后：
       python3 inject_notes.py "FY26 EO Training v3_FIXED-5.pptx"
     不传文件名则默认找当前目录下的 "FY26 EO Training v3_FIXED-5.pptx"。
  3) 生成新文件： "FY26 EO Training v3_FIXED-5_带备注.pptx"（原文件不动）。

说明：
  - 按「页码 = 幻灯片序号」一一对应注入；若某页 PPT 里没有备注占位符，会自动创建。
  - 只改备注，不动任何正文/图片/视频，保真。
"""

import os
import re
import sys


def load_notes(md_path):
    """从 markdown 解析出 {页码: 备注全文}。"""
    with open(md_path, encoding="utf-8") as f:
        text = f.read()

    notes = {}
    # 以 "## 第 N 页" 切分
    parts = re.split(r"^##\s+第\s*(\d+)\s*页", text, flags=re.MULTILINE)
    # parts: [前言, '1', body1, '2', body2, ...]
    for i in range(1, len(parts) - 1, 2):
        page = int(parts[i])
        body = parts[i + 1]
        # 去掉结尾的章节分隔线及之后内容（"### 📌 交付与版本说明" 等）
        body = re.split(r"^###\s", body, flags=re.MULTILINE)[0]
        # 标题行（· 后面的内容）拼回备注开头，保留上下文
        lines = body.splitlines()
        title = lines[0].lstrip("· ").strip() if lines else ""
        rest = "\n".join(lines[1:]) if len(lines) > 1 else ""
        # 清掉 markdown 粗体/分隔符号，让备注框里更干净
        cleaned = rest.replace("**", "").strip().strip("-").strip()
        cleaned = re.sub(r"\n-{3,}\s*$", "", cleaned).strip()
        note = (("【" + title + "】\n\n") if title else "") + cleaned
        notes[page] = note.strip()
    return notes


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    md_path = os.path.join(here, "FY26_EO_Training_备注讲解.md")
    if not os.path.exists(md_path):
        sys.exit("找不到 FY26_EO_Training_备注讲解.md，请把它和本脚本放同一目录。")

    pptx_path = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        here, "FY26 EO Training v3_FIXED-5.pptx")
    if not os.path.exists(pptx_path):
        sys.exit(f"找不到 PPTX：{pptx_path}\n请把文件名作为参数传入，或放到本脚本同目录。")

    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("缺少依赖，请先运行： pip3 install python-pptx")

    notes = load_notes(md_path)
    print(f"已从 markdown 解析出 {len(notes)} 页备注。")

    prs = Presentation(pptx_path)
    n_slides = len(prs.slides)
    print(f"PPTX 共 {n_slides} 页。")

    written = 0
    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in notes:
            print(f"  · 第 {idx} 页：markdown 里没有对应备注，跳过。")
            continue
        slide.notes_slide.notes_text_frame.text = notes[idx]
        written += 1

    if n_slides != len(notes):
        print(f"⚠️  注意：PPT 页数({n_slides}) 与备注页数({len(notes)}) 不一致，"
              f"已按序号尽量对齐，请抽查首尾几页。")

    base, ext = os.path.splitext(pptx_path)
    out_path = base + "_带备注" + ext
    prs.save(out_path)
    print(f"✅ 完成！已写入 {written} 页备注。")
    print(f"成品： {out_path}")
    print("原文件未改动，可直接用成品文件演示。")


if __name__ == "__main__":
    main()
